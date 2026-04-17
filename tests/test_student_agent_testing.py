from fastapi.testclient import TestClient
from pptx import Presentation

from backend.app import app
from backend.aggregation import aggregate_runs_for_deck
from backend.models import LectureDeck, QAQuestion, SlideBlock, TaskType
from backend.profiles import DEFAULT_PROFILES
from backend.runner import run_evaluation_job, select_profiles, select_tasks
from backend.storage import DECK_STORE, JOB_STORE, RUN_STORE, SIGNAL_STORE
from backend.tasks import DEFAULT_TASKS


def reset_stores() -> None:
    DECK_STORE.clear()
    JOB_STORE.clear()
    RUN_STORE.clear()
    SIGNAL_STORE.clear()


def build_demo_deck() -> LectureDeck:
    return LectureDeck(
        deck_id="deck_test_001",
        module_tag="intro_ai",
        week_tag="week_1",
        topic_tags=["machine learning basics", "supervised learning"],
        slides=[
            SlideBlock(
                slide_id="s1",
                title="What is Machine Learning?",
                key_points=["ML learns patterns from data"],
                terms=["machine learning", "data", "pattern"],
                examples=["spam detection"],
                raw_text="Machine learning is a method for learning patterns from data.",
            ),
            SlideBlock(
                slide_id="s2",
                title="Supervised Learning",
                key_points=["Uses labelled examples", "Predicts outputs from inputs"],
                terms=["label", "training data", "prediction"],
                examples=[],
                raw_text="Supervised learning maps inputs to outputs using labelled training examples.",
            ),
        ],
    )


def build_demo_payload() -> dict:
    return {
        "deck": build_demo_deck().model_dump(),
        "profile_names": ["weak", "average", "strong"],
        "task_types": ["explain_back", "short_qa", "confusion_report", "coverage_check"],
        "auto_aggregate": True,
        "use_mock": True,
    }


def build_demo_questions() -> list[QAQuestion]:
    return [
        QAQuestion(
            question_id=f"q{index}",
            prompt=f"What is concept {index}?",
            correct_answer=f"Correct answer {index}",
            slide_refs=["s1" if index <= 5 else "s2"],
            concept_tags=[f"concept_{index}"],
        )
        for index in range(1, 11)
    ]


def test_mock_pipeline_generates_summary() -> None:
    reset_stores()
    deck = build_demo_deck()
    profiles = select_profiles([profile.profile_name for profile in DEFAULT_PROFILES], DEFAULT_PROFILES)
    tasks = select_tasks([task.task_type for task in DEFAULT_TASKS], DEFAULT_TASKS)

    job = run_evaluation_job(deck=deck, profiles=profiles, tasks=tasks, use_mock=True)
    summary = aggregate_runs_for_deck(deck.deck_id)

    assert job.status.value == "completed"
    assert len(job.run_ids) == len(deck.slides) * 10
    assert summary.tested_deck_id == deck.deck_id
    assert summary.weak_topic_counts
    assert summary.repeated_confusion_points
    assert summary.evidence_refs
    assert summary.governance_trace is not None
    assert summary.governance_trace.run_ids


def test_coverage_check_runs_only_for_strong_profile() -> None:
    reset_stores()
    deck = build_demo_deck()
    profiles = select_profiles([profile.profile_name for profile in DEFAULT_PROFILES], DEFAULT_PROFILES)
    tasks = select_tasks([TaskType.coverage_check], DEFAULT_TASKS)

    job = run_evaluation_job(deck=deck, profiles=profiles, tasks=tasks, use_mock=True)
    runs = [RUN_STORE[run_id] for run_id in job.run_ids]

    assert runs
    assert all(run.student_profile.value == "strong" for run in runs)


def test_profiles_produce_distinguishable_patterns() -> None:
    reset_stores()
    deck = build_demo_deck()
    profiles = select_profiles([profile.profile_name for profile in DEFAULT_PROFILES], DEFAULT_PROFILES)
    tasks = select_tasks([TaskType.explain_back], DEFAULT_TASKS)

    job = run_evaluation_job(deck=deck, profiles=profiles, tasks=tasks, use_mock=True)
    runs = [RUN_STORE[run_id] for run_id in job.run_ids]
    by_profile = {run.student_profile.value: run.answer for run in runs if run.slide_refs == ["s1"]}

    assert "terminology_overload" in by_profile["weak"].confusion_tags
    assert "pace_jump" in by_profile["average"].confusion_tags
    assert "oversimplified_explanation" in by_profile["strong"].misconception_flags


def test_run_endpoint_returns_job_and_summary() -> None:
    reset_stores()
    client = TestClient(app)

    response = client.post("/evaluation/run", json=build_demo_payload())

    assert response.status_code == 200
    payload = response.json()
    assert payload["job"]["status"] == "completed"
    assert payload["summary"]["tested_deck_id"] == "deck_test_001"
    assert len(payload["job"]["run_ids"]) == 20
    assert payload["summary"]["representative_evidence"]


def test_qa_improvement_endpoint_returns_uplift_and_handoff() -> None:
    reset_stores()
    client = TestClient(app)
    response = client.post(
        "/evaluation/qa-improvement",
        json={
            "deck": build_demo_deck().model_dump(),
            "questions": [question.model_dump() for question in build_demo_questions()],
            "profile_name": "weak",
            "qa_support_content": [
                "Explain the label concept with one worked example.",
                "Clarify the difference between data, model, and pattern.",
            ],
            "target_collaborator": "content_revision_collaborator",
            "use_mock": True,
        },
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["before_pass"]["correct_count"] == 2
    assert payload["after_pass"]["correct_count"] == 8
    assert payload["improvement_delta"] == 6
    assert payload["model_used"] == "mock_student_response"
    assert payload["collaborator_handoff"]["unresolved_count"] == 2
    assert len(payload["collaborator_handoff"]["items"]) == 2
    assert payload["collaborator_handoff"]["schema_version"] == "collaborator_handoff.v1"
    assert payload["collaborator_handoff"]["handoff_type"] == "content_revision_handoff"
    assert payload["collaborator_handoff"]["source_system"] == "student_agent_testing"
    assert payload["collaborator_handoff"]["target_collaborator"] == "content_revision_collaborator"
    assert payload["collaborator_handoff"]["items"][0]["target_type"] == "question"
    assert payload["collaborator_handoff"]["items"][0]["issue_types"] == [
        "incorrect_after_support",
        "content_revision_candidate",
    ]
    assert payload["collaborator_handoff"]["items"][0]["affected_profiles"] == ["weak"]
    assert payload["collaborator_handoff"]["items"][0]["priority"] == "high"
    assert payload["collaborator_handoff"]["items"][0]["student_answer"]
    assert payload["collaborator_handoff"]["items"][0]["correct_answer"]
    assert payload["collaborator_handoff"]["items"][0]["missed_concepts"]


def test_pptx_ingest_returns_structured_deck(tmp_path) -> None:
    reset_stores()
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Intro to Supervised Learning"
    slide.placeholders[1].text = "Labels map inputs to outputs\nExample: spam filtering"

    pptx_path = tmp_path / "sample_deck.pptx"
    presentation.save(str(pptx_path))

    client = TestClient(app)
    response = client.post(
        "/deck/ingest/pptx",
        json={
            "pptx_path": str(pptx_path),
            "deck_id": "deck_from_pptx",
            "module_tag": "intro_ai",
            "week_tag": "week_2",
            "topic_tags": ["supervised learning"],
        },
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["deck"]["deck_id"] == "deck_from_pptx"
    assert payload["deck"]["slides"][0]["title"] == "Intro to Supervised Learning"
    assert "Labels map inputs to outputs" in payload["deck"]["slides"][0]["raw_text"]


def test_job_status_includes_answer_text_and_evidence_refs() -> None:
    reset_stores()
    client = TestClient(app)
    run_response = client.post("/evaluation/run", json=build_demo_payload())
    job_id = run_response.json()["job"]["job_id"]

    status_response = client.get(f"/evaluation/{job_id}")

    assert status_response.status_code == 200
    payload = status_response.json()
    assert payload["runs"][0]["answer_text"]
    assert payload["runs"][0]["evidence_refs"]


def test_summary_includes_governance_trace() -> None:
    reset_stores()
    client = TestClient(app)
    client.post("/evaluation/run", json=build_demo_payload())

    response = client.get("/evaluation/summary/deck_test_001")

    assert response.status_code == 200
    payload = response.json()
    assert payload["governance_trace"]["signal_id"] == "signal_deck_test_001"
    assert payload["governance_trace"]["run_ids"]
    assert payload["governance_trace"]["evidence_refs"]
