from __future__ import annotations

from pathlib import Path

from backend.models import LectureDeck, SlideBlock


def pptx_to_lecture_deck(
    pptx_path: str,
    deck_id: str,
    module_tag: str,
    week_tag: str,
    topic_tags: list[str],
) -> LectureDeck:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is not installed.") from exc

    path = Path(pptx_path)
    if not path.exists():
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

    presentation = Presentation(str(path))
    slides: list[SlideBlock] = []

    for index, slide in enumerate(presentation.slides, start=1):
        text_items = _extract_slide_text(slide)
        title = text_items[0] if text_items else f"Slide {index}"
        body_items = text_items[1:] if len(text_items) > 1 else []
        raw_text = "\n".join(text_items).strip() or None
        key_points = body_items[:5]
        terms = _extract_terms(title, body_items)
        examples = [item for item in body_items if "example" in item.lower() or "e.g." in item.lower()]

        slides.append(
            SlideBlock(
                slide_id=f"s{index}",
                title=title,
                key_points=key_points,
                terms=terms[:6],
                examples=examples[:3],
                raw_text=raw_text,
            )
        )

    return LectureDeck(
        deck_id=deck_id,
        module_tag=module_tag,
        week_tag=week_tag,
        topic_tags=topic_tags,
        slides=slides,
    )


def _extract_slide_text(slide) -> list[str]:
    items: list[str] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = "\n".join(paragraph.text.strip() for paragraph in shape.text_frame.paragraphs if paragraph.text.strip()).strip()
        if text:
            items.extend([line.strip() for line in text.splitlines() if line.strip()])
    return items


def _extract_terms(title: str, body_items: list[str]) -> list[str]:
    candidates: list[str] = []
    for text in [title] + body_items:
        parts = [part.strip(" ,.:;()") for part in text.replace("/", " ").split()]
        for part in parts:
            if len(part) >= 4 and part[0].isalnum():
                lowered = part.lower()
                if lowered not in {item.lower() for item in candidates}:
                    candidates.append(part)
    return candidates
